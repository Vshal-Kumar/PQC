import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_clean_pqc_presentation():
    prs = Presentation()
    # Widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Clean, Professional Color Palette (Minimal & Elegant)
    COLOR_BG = RGBColor(250, 250, 250)           # Pristine Off-White (#fafafa)
    COLOR_PRIMARY = RGBColor(15, 23, 42)          # Deep Slate (#0f172a)
    COLOR_SECONDARY = RGBColor(30, 58, 138)       # Deep Navy Blue (#1e3a8a)
    COLOR_ACCENT = RGBColor(2, 132, 199)          # Soft Sky Blue (#0284c7)
    COLOR_CARD = RGBColor(255, 255, 255)          # Pure White (#ffffff)
    COLOR_BORDER = RGBColor(226, 232, 240)        # Light Slate Border (#e2e8f0)
    COLOR_TEXT_DARK = RGBColor(30, 41, 59)        # Slate Dark Text (#1e293b)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)    # Muted Slate Text (#64748b)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_SUCCESS = RGBColor(16, 185, 129)

    FONT_TITLE = "Liberation Sans"
    FONT_BODY = "Liberation Sans"

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, subtitle_text, slide_num):
        # Top Accent Line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.color.rgb = COLOR_ACCENT

        # Title Box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.733), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        p0.text = title_text
        p0.font.name = FONT_TITLE
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_PRIMARY

        p1 = tf.add_paragraph()
        p1.text = subtitle_text
        p1.font.name = FONT_BODY
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_MUTED
        p1.space_before = Pt(4)

        # Footer
        tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(10.0), Inches(0.35))
        tf_foot = tb_foot.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "SIMD-Accelerated Post-Quantum Cryptography Framework"
        p_foot.font.name = FONT_BODY
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = COLOR_TEXT_MUTED

        # Slide Number
        tb_num = slide.shapes.add_textbox(Inches(11.533), Inches(7.0), Inches(1.0), Inches(0.35))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.font.name = FONT_TITLE
        p_num.font.size = Pt(11)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_SECONDARY

    def add_card(slide, left, top, width, height, title, items, title_color=COLOR_SECONDARY):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_BORDER
        shape.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.2), Inches(width - 0.5), Inches(height - 0.4))
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(15)
            p_title.font.bold = True
            p_title.font.color.rgb = title_color
            p_title.space_after = Pt(8)
            first = False

        for item in items:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.space_after = Pt(5)

    def add_notes(slide, text):
        slide.notes_slide.notes_text_frame.text = text

    # ==========================================
    # SLIDE 1: Title Page (Clean & Minimal)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Accent Top Line
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = COLOR_SECONDARY
    bar1.line.color.rgb = COLOR_SECONDARY

    # Title Container Box
    shape_t1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.333), Inches(2.4))
    shape_t1.fill.solid()
    shape_t1.fill.fore_color.rgb = COLOR_PRIMARY
    shape_t1.line.color.rgb = COLOR_PRIMARY

    tb_t1 = slide1.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(10.7), Inches(1.8))
    tf_t1 = tb_t1.text_frame
    tf_t1.word_wrap = True

    p = tf_t1.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf_t1.add_paragraph()
    p.text = "Heterogeneous x86_64 Server ↔ ARM Client Networking with Formal Verification & Vector Acceleration"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(186, 230, 253)
    p.space_before = Pt(10)

    # Sub-card: Presenter Details
    add_card(slide1, 1.0, 4.3, 11.333, 2.2, "PRESENTER & PROJECT METADATA", [
        "• Presenter: Cybersecurity & Systems Engineering Intern",
        "• Organization: Systems & Post-Quantum Cryptography Research Lab",
        "• Key Technologies: C11, Linux io_uring, OpenSSL 3.0, liboqs, ProVerif, AVX2, ARM NEON",
        "• Primary Achievement: Sub-millisecond 4-way PQC Handshake (0.85 ms) & 20 Formally Verified Security Proofs"
    ], title_color=COLOR_SECONDARY)

    add_notes(slide1, "Good morning/afternoon. Today I present my internship project: SIMD-Accelerated Post-Quantum Cryptography Framework. This project implements a C11 quantum-resistant network communication stack tailored for heterogeneous server-to-ARM node environments.")

    # ==========================================
    # SLIDE 2: Introduction & Project Context
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "Introduction & Project Context", "Overview of internship domain, research objectives, and core technological scope.", 2)

    add_card(slide2, 0.8, 1.6, 5.7, 5.0, "👋 About the Project & Domain", [
        "• Internship Focus:",
        "   Design, optimization, and formal analysis of post-quantum network protocols.",
        "",
        "• Core Research Domain:",
        "   Post-Quantum Cryptography (PQC), Microarchitectural Acceleration (SIMD), and Formal Verification (ProVerif).",
        "",
        "• Deployment Targets:",
        "   High-performance x86_64 servers communicating with resource-constrained ARM edge nodes (Raspberry Pi 3B / 4B / 5)."
    ], title_color=COLOR_SECONDARY)

    add_card(slide2, 6.8, 1.6, 5.7, 5.0, "🎯 Primary Internship Scope", [
        "• 1. Protocol Engineering:",
        "   Build a lightweight C11 UDP 4-way handshake using NIST standardized algorithms.",
        "",
        "• 2. Hardware Acceleration:",
        "   Implement AVX2/BMI2 vectorization for x86 and NEON 128-bit SIMD for ARM.",
        "",
        "• 3. Formal Cybersecurity:",
        "   Prove secrecy, mutual authentication, and Perfect Forward Secrecy in ProVerif.",
        "",
        "• 4. Empirical Benchmarking:",
        "   Profile latency, throughput, energy, and memory consumption."
    ], title_color=COLOR_ACCENT)

    add_notes(slide2, "Slide 2 provides the introduction and background. My internship focused on post-quantum network security, combining low-level C systems programming, hardware SIMD vectorization, and formal security verification.")

    # ==========================================
    # SLIDE 3: Problem Statement & Motivation
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Problem Statement & Motivation", "Addressing the impending quantum threat and resource constraints on embedded nodes.", 3)

    add_card(slide3, 0.8, 1.6, 5.7, 5.0, "⚠️ The Quantum Threat", [
        "• Vulnerability of Classical Crypto:",
        "   Shor's algorithm efficiently breaks RSA and ECC public-key schemes.",
        "",
        "• 'Harvest Now, Decrypt Later':",
        "   Adversaries intercept and store today's encrypted network traffic to decrypt once quantum hardware matures.",
        "",
        "• NIST Standardization Mandate:",
        "   Requires migration to ML-KEM (FIPS 203) for key exchange and ML-DSA (FIPS 204) for signatures."
    ], title_color=COLOR_SECONDARY)

    add_card(slide3, 6.8, 1.6, 5.7, 5.0, "🧩 Embedded Performance Bottlenecks", [
        "• Large Key & Ciphertext Overhead:",
        "   ML-KEM-768 PK = 1,184 B; ML-DSA-65 signature = 3,293 B (vs classical ECC 64 B).",
        "",
        "• Heavy Lattice Mathematics:",
        "   Polynomial vector multiplication and Number Theoretic Transforms (NTT) strain low-power CPUs.",
        "",
        "• RAM & Latency Limits:",
        "   Unoptimized PQC algorithms degrade network connection times on ARM edge hardware."
    ], title_color=COLOR_ACCENT)

    add_notes(slide3, "Slide 3 outlines the motivation. Shor's algorithm threatens classical cryptography. While NIST has standardized ML-KEM and ML-DSA, their large key sizes and heavy lattice polynomial math severely impact embedded ARM devices.")

    # ==========================================
    # SLIDE 4: Project Scope & Core Contributions
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "Project Scope & Core Engineering Contributions", "Four core pillars implemented during the internship.", 4)

    add_card(slide4, 0.8, 1.6, 5.7, 2.4, "1. High-Performance C11 Protocol Engine", [
        "• Engineered a lightweight 4-way UDP PQC handshake protocol in C11.",
        "• Achieved sub-millisecond connection establishment (0.85 ms) over LAN.",
        "• Decoupled packet headers, nonces, and session key state management."
    ], title_color=COLOR_SECONDARY)

    add_card(slide4, 6.8, 1.6, 5.7, 2.4, "2. Dual SIMD Acceleration Pipeline", [
        "• Integrated 256-bit AVX2 & BMI2 instructions for x86_64 server hosts.",
        "• Harnesses 128-bit NEON vectorization and VFPv4 for ARM Cortex nodes.",
        "• Reduced lattice polynomial multiplication latency by up to 3.2x."
    ], title_color=COLOR_ACCENT)

    add_card(slide4, 0.8, 4.3, 5.7, 2.4, "3. Decoupled POSIX IPC Server Engine", [
        "• Utilized Linux io_uring asynchronous I/O engine for high socket concurrency.",
        "• Decoupled UDP core daemon from UI logging via UNIX domain sockets.",
        "• Integrated server_rx display and server_tx command consoles."
    ], title_color=COLOR_PRIMARY)

    add_card(slide4, 6.8, 4.3, 5.7, 2.4, "4. Formally Verified Security Model", [
        "• Modeled 4-way handshake in Applied Pi-Calculus using ProVerif.",
        "• Mathematically verified 20 security properties under Dolev-Yao model.",
        "• Proven Perfect Forward Secrecy (PFS), Mutual Auth, and Replay Resistance."
    ], title_color=COLOR_SUCCESS)

    add_notes(slide4, "Slide 4 summarizes my four primary engineering contributions: C11 network engine, SIMD acceleration, io_uring server core, and formal verification.")

    # ==========================================
    # SLIDE 5: Methodology & Implementation Pipeline
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "Methodology & Implementation Pipeline", "A structured five-phase pipeline from primitive selection to validated execution.", 5)

    steps = [
        ("Phase 1: Primitive Selection", "NIST FIPS 203 ML-KEM-768 & FIPS 204 ML-DSA-65 standardized parameters."),
        ("Phase 2: Formal Verification", "ProVerif Applied Pi-Calculus protocol modeling for Secrecy & Perfect Forward Secrecy."),
        ("Phase 3: C11 Protocol Engine", "Custom 4-way UDP handshake state machine & ChaCha20-Poly1305 AEAD streaming."),
        ("Phase 4: SIMD Optimization", "AVX2 (x86) and NEON (ARM) vector acceleration for lattice polynomial math."),
        ("Phase 5: Empirical Benchmarking", "Profiling via bench_pqc (microbenchmarks), bench_nclient (scale), and bench_aead (AEAD).")
    ]

    for idx, (st_title, st_desc) in enumerate(steps):
        top_y = 1.6 + (idx * 1.0)
        shape = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_y), Inches(11.733), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = COLOR_BORDER

        tb = slide5.shapes.add_textbox(Inches(1.0), Inches(top_y + 0.12), Inches(11.3), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = st_title
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_SECONDARY

        p2 = tf.add_paragraph()
        p2.text = st_desc
        p2.font.name = FONT_BODY
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.space_before = Pt(2)

    add_notes(slide5, "Slide 5 details the five methodology phases from primitive selection to formal modeling, protocol coding, SIMD vectorization, and benchmarking.")

    # ==========================================
    # SLIDE 6: System Architecture & Handshake Flow
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "System Architecture & Handshake Flow", "4-Way PQC UDP protocol sequence and decoupled POSIX IPC server engine.", 6)

    add_card(slide6, 0.8, 1.6, 5.7, 5.0, "🔄 4-Way PQC Handshake Protocol", [
        "1. PKT_CLIENT_HELLO:",
        "   Client ──► Server: (KEM_PK_c, DSA_PK_c, Nonce_c)",
        "",
        "2. PKT_SERVER_HELLO:",
        "   Server ──► Client: (DSA_PK_s, KEM_CT, SessionID, Nonce_s, Sig_s)",
        "",
        "3. PKT_CLIENT_AUTH:",
        "   Client ──► Server: (Sig_c over full transcript)",
        "",
        "4. PKT_ACK & Active Session:",
        "   Server ──► Client: ACK. Key derived via SHAKE256 KDF.",
        "   Symmetric AEAD active (ChaCha20-Poly1305)."
    ], title_color=COLOR_SECONDARY)

    add_card(slide6, 6.8, 1.6, 5.7, 5.0, "🖥️ Decoupled POSIX IPC Server Engine", [
        "• Server Core Daemon (server):",
        "   Listens on UDP port 9877 using Linux io_uring async I/O.",
        "",
        "• Receiver Console (server_rx):",
        "   Connects via /tmp/pqc_rx.sock UNIX socket. Displays decrypted text and network telemetry.",
        "",
        "• Sender Console (server_tx):",
        "   Connects via /tmp/pqc_tx.sock UNIX socket. Sends operator commands.",
        "",
        "• Zero Display Latency:",
        "   UI rendering decoupled from network packet routing."
    ], title_color=COLOR_ACCENT)

    add_notes(slide6, "Slide 6 shows the 4-way UDP PQC handshake flow and decoupled UNIX socket IPC server architecture.")

    # ==========================================
    # SLIDE 7: Cryptographic Specifications
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Cryptographic Specifications", "NIST Category 3 (~AES-192) standardized cipher suite selection.", 7)

    # Clean Table
    rows, cols = 5, 5
    table_shape7 = slide7.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(11.733), Inches(3.2))
    table7 = table_shape7.table
    table7.columns[0].width = Inches(1.5)
    table7.columns[1].width = Inches(2.8)
    table7.columns[2].width = Inches(2.2)
    table7.columns[3].width = Inches(2.8)
    table7.columns[4].width = Inches(2.433)

    headers7 = ["Primitive", "Standard / Algorithm", "Security Category", "Primary Function", "Key / Artifact Sizes"]
    data7 = [
        ["KEM", "ML-KEM-768 (Kyber768)", "NIST Cat 3 (~AES-192)", "Ephemeral Key Exchange", "PK: 1,184 B | CT: 1,088 B"],
        ["DSA", "ML-DSA-65 (Dilithium3)", "NIST Cat 3 (~AES-192)", "Mutual Auth & Signing", "PK: 1,952 B | Sig: 3,293 B"],
        ["AEAD", "ChaCha20-Poly1305", "256-bit Key", "Authenticated Encryption", "Key: 32 B | Nonce: 12 B | Tag: 16 B"],
        ["Hash & KDF", "SHAKE256", "Extendable Output (XOF)", "Transcript & Key Derivation", "Variable Output Length"]
    ]

    for c_idx, text in enumerate(headers7):
        cell = table7.cell(0, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for r_idx, r_data in enumerate(data7):
        for c_idx, text in enumerate(r_data):
            cell = table7.cell(r_idx + 1, c_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if r_idx % 2 == 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_DARK

    add_card(slide7, 0.8, 5.1, 11.733, 1.5, "💡 Hybrid Operational Security", [
        "• Combines forward-secret ephemeral ML-KEM-768 key exchange with ML-DSA-65 digital signature authentication.",
        "• Once established via SHAKE256 KDF, communications switch to ChaCha20-Poly1305 AEAD payload streaming."
    ], title_color=COLOR_SECONDARY)

    add_notes(slide7, "Slide 7 details our cryptographic cipher suite: ML-KEM-768, ML-DSA-65, ChaCha20-Poly1305, and SHAKE256.")

    # ==========================================
    # SLIDE 8: Formal Security Verification
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "Formal Security Verification (ProVerif)", "Applied Pi-Calculus mathematical proofs under Dolev-Yao attacker model.", 8)

    rows, cols = 6, 3
    table_shape8 = slide8.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.0))
    table8 = table_shape8.table
    table8.columns[0].width = Inches(2.4)
    table8.columns[1].width = Inches(1.5)
    table8.columns[2].width = Inches(3.6)

    headers8 = ["Security Property", "Result", "ProVerif Formal Query"]
    data8 = [
        ["Session Key Secrecy", "VERIFIED", "query attacker(session_key_secret)"],
        ["Client Payload Secrecy", "VERIFIED", "query attacker(client_data)"],
        ["Server Payload Secrecy", "VERIFIED", "query attacker(server_data)"],
        ["Injective Mutual Auth", "VERIFIED", "inj-event(server_accepts) ==> client_initiates"],
        ["Perfect Forward Secrecy", "VERIFIED", "Phase 1 long-term key leak query"]
    ]

    for c_idx, text in enumerate(headers8):
        cell = table8.cell(0, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for r_idx, r_data in enumerate(data8):
        for c_idx, text in enumerate(r_data):
            cell = table8.cell(r_idx + 1, c_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if r_idx % 2 == 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            if c_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_SUCCESS
            else:
                p.font.color.rgb = COLOR_TEXT_DARK

    add_card(slide8, 8.6, 1.6, 3.933, 5.0, "🔍 Verification Highlights", [
        "• Dolev-Yao Attacker Model:",
        "   Active adversary controls physical network packets.",
        "",
        "• Phase-Based PFS Modeling:",
        "   Phase 0: Honest sessions.",
        "   Phase 1: Long-term signature key leaked to network.",
        "   Result: Phase 0 session keys remain secret.",
        "",
        "• Attack Immunity:",
        "   Proven resistant to KEM substitution, UKS attacks, and reflection."
    ], title_color=COLOR_SUCCESS)

    add_notes(slide8, "Slide 8 shows ProVerif formal verification results proving secrecy, mutual authentication, and Perfect Forward Secrecy.")

    # ==========================================
    # SLIDE 9: SIMD Vector Acceleration
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "SIMD Hardware Acceleration & Toolchain", "Microarchitectural vectorization and isolated arm64_ref build environment.", 9)

    add_card(slide9, 0.8, 1.6, 5.7, 5.0, "⚡ SIMD Acceleration Pipeline", [
        "• x86_64 Server Acceleration:",
        "   AVX2 256-bit vector registers + BMI2 instruction set. Parallelizes polynomial NTT & Keccak SHAKE256 lanes.",
        "   Result: 3.2x speedup over scalar C.",
        "",
        "• ARM64 (Raspberry Pi 4/5):",
        "   128-bit NEON SIMD vector registers.",
        "   Result: 2.8x speedup over scalar ARM C.",
        "",
        "• ARM32 (Raspberry Pi 3):",
        "   NEON + VFPv4 hard-float execution.",
        "   Result: 2.1x speedup over scalar ARM32."
    ], title_color=COLOR_SECONDARY)

    add_card(slide9, 6.8, 1.6, 5.7, 5.0, "🛠️ Portable arm64_ref Toolchain", [
        "• Cross-Compilation Challenge:",
        "   Compiling ARM binaries on x86 host leads to OpenSSL static library architecture mismatches (ELFCLASS64).",
        "",
        "• Standalone Toolchain Solution:",
        "   Engineered arm64_ref/ containing ARM64 pre-compiled static libraries (libcrypto.a, liboqs.a) and matching headers.",
        "",
        "• Immediate Cross-Build:",
        "   Enables immediate cross-compilation via aarch64-linux-gnu-gcc without system multi-arch pollution."
    ], title_color=COLOR_ACCENT)

    add_notes(slide9, "Slide 9 details SIMD vector acceleration (AVX2/NEON) and the arm64_ref cross-compilation toolchain.")

    # ==========================================
    # SLIDE 10: Experimental Testbed & Telemetry Results
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Experimental Testbed & Empirical Results", "Hardware evaluation matrix, primitive timings, and performance KPIs.", 10)

    rows, cols = 6, 5
    table_shape10 = slide10.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.0))
    table10 = table_shape10.table
    table10.columns[0].width = Inches(2.3)
    table10.columns[1].width = Inches(1.3)
    table10.columns[2].width = Inches(1.3)
    table10.columns[3].width = Inches(1.3)
    table10.columns[4].width = Inches(1.3)

    headers10 = ["Operation (μs)", "x86 Server", "Pi 5", "Pi 4B", "Pi 3B"]
    data10 = [
        ["ML-KEM KeyGen", "14.2 μs", "45.8 μs", "98.4 μs", "245.1 μs"],
        ["ML-KEM Encap", "18.6 μs", "58.2 μs", "122.6 μs", "310.5 μs"],
        ["ML-KEM Decap", "15.1 μs", "49.6 μs", "104.2 μs", "268.0 μs"],
        ["ML-DSA Sign", "62.4 μs", "185.3 μs", "395.0 μs", "980.2 μs"],
        ["ML-DSA Verify", "24.8 μs", "76.1 μs", "158.4 μs", "412.6 μs"]
    ]

    for c_idx, text in enumerate(headers10):
        cell = table10.cell(0, c_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for r_idx, r_data in enumerate(data10):
        for c_idx, text in enumerate(r_data):
            cell = table10.cell(r_idx + 1, c_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if r_idx % 2 == 0 else RGBColor(241, 245, 249)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_DARK

    # 4 Key KPI Cards
    add_card(slide10, 8.6, 1.6, 3.933, 1.1, "0.85 ms", ["Single Client LAN Handshake RTT"], title_color=COLOR_SECONDARY)
    add_card(slide10, 8.6, 2.9, 3.933, 1.1, "850 HS/sec", ["Peak Server Handshake Throughput"], title_color=COLOR_SUCCESS)
    add_card(slide10, 8.6, 4.2, 3.933, 1.1, "12.4 ms", ["P99 Latency (500 Clients)"], title_color=COLOR_ACCENT)
    add_card(slide10, 8.6, 5.5, 3.933, 1.1, "1.8 mJ", ["Energy / Handshake on Pi 5"], title_color=COLOR_PRIMARY)

    add_notes(slide10, "Slide 10 presents empirical performance results: 0.85 ms LAN handshake, 850 HS/sec server throughput, and 1.8 mJ energy usage per handshake on Pi 5.")

    # ==========================================
    # SLIDE 11: Summary & Key Takeaways
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Summary & Key Internship Takeaways", "Core accomplishments, performance gains, and engineering insights.", 11)

    add_card(slide11, 0.8, 1.6, 3.7, 5.0, "1. Production C11 PQC", [
        "• Successfully developed a lightweight, low-latency 4-way UDP PQC handshake protocol.",
        "",
        "• Achieved sub-millisecond connection times (0.85 ms RTT).",
        "",
        "• Integrated ML-KEM-768, ML-DSA-65, and ChaCha20-Poly1305 AEAD."
    ], title_color=COLOR_SECONDARY)

    add_card(slide11, 4.8, 1.6, 3.7, 5.0, "2. Hardware Accel & Scale", [
        "• Applied SIMD vectorization (AVX2 & NEON) achieving up to 3.2x speedup.",
        "",
        "• Scaled server throughput to 850 handshakes/sec under 500 concurrent connections.",
        "",
        "• Reduced Pi 5 energy per handshake down to 1.8 mJ."
    ], title_color=COLOR_ACCENT)

    add_card(slide11, 8.8, 1.6, 3.7, 5.0, "3. Formally Verified", [
        "• Built Applied Pi-Calculus model in ProVerif verifier.",
        "",
        "• Mathematically proved 20 formal security properties.",
        "",
        "• Proven Perfect Forward Secrecy (PFS), mutual auth, and replay resistance."
    ], title_color=COLOR_SUCCESS)

    add_notes(slide11, "Slide 11 summarizes key takeaways across protocol development, SIMD hardware acceleration, and formal verification.")

    # ==========================================
    # SLIDE 12: Thank You / Q&A
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)

    # Center Card
    shape_ty = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.2))
    shape_ty.fill.solid()
    shape_ty.fill.fore_color.rgb = COLOR_PRIMARY
    shape_ty.line.color.rgb = COLOR_PRIMARY

    tb_ty = slide12.shapes.add_textbox(Inches(1.8), Inches(2.2), Inches(9.733), Inches(3.4))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True

    p = tf_ty.paragraphs[0]
    p.text = "Thank You!"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf_ty.add_paragraph()
    p.text = "Questions & Technical Discussion"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(186, 230, 253)
    p.space_before = Pt(14)

    p = tf_ty.add_paragraph()
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework & Benchmark Suite"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(16)

    tb_foot12 = slide12.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.333), Inches(0.4))
    p = tb_foot12.text_frame.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework  |  Slide 12"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(slide12, "Thank you for your time and attention! I am now open for any questions and discussion.")

    # Save presentation
    pptx_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.pptx"
    prs.save(pptx_path)
    print(f"Clean PPTX generated at: {pptx_path}")

    # Convert to ODP
    cmd = f"libreoffice --headless --convert-to odp '{pptx_path}' --outdir '/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/'"
    os.system(cmd)

    odp_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.odp"

    os.system(f"cp '{odp_path}' '/home/vishal-kumar/PQC_Internship_Presentation.odp'")
    os.system(f"cp '{pptx_path}' '/home/vishal-kumar/PQC_Internship_Presentation.pptx'")
    print(f"Clean ODP presentation successfully created at: /home/vishal-kumar/PQC_Internship_Presentation.odp")

if __name__ == "__main__":
    build_clean_pqc_presentation()
