import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_refined_pqc_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # Color Palette matching Reference Deck (Quantum_Trajectory_Deck)
    COLOR_BG = RGBColor(248, 250, 252)            # Clean Light Slate (#f8fafc)
    COLOR_HEADER_BG = RGBColor(11, 37, 69)        # Deep Space Navy (#0b2545)
    COLOR_PILL_BG = RGBColor(11, 61, 145)         # Royal Deep Blue (#0b3d91)
    COLOR_ACCENT_TEAL = RGBColor(28, 114, 147)    # Teal Cyan (#1c7293)
    COLOR_CARD_BG = RGBColor(255, 255, 255)       # Pure White Card (#ffffff)
    COLOR_CARD_BORDER = RGBColor(211, 222, 235)   # Soft Border (#d3deeb)
    COLOR_TEXT_MAIN = RGBColor(15, 23, 42)        # Deep Slate Text (#0f172a)
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)    # Muted Text (#64748b)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_EMERALD = RGBColor(16, 185, 129)
    COLOR_AMBER = RGBColor(217, 119, 6)

    FONT_TITLE = "Inter"
    FONT_BODY = "Inter"
    FONT_MONO = "DejaVu Sans Mono"

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, category_tag, title_text, subtitle_text, slide_num):
        # Top Accent Cyan Divider Line
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = COLOR_ACCENT_TEAL
        bar.line.color.rgb = COLOR_ACCENT_TEAL

        # Category Tag (Pill Shape)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.35), Inches(2.2), Inches(0.35))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COLOR_PILL_BG
        pill.line.color.rgb = COLOR_PILL_BG
        tf_pill = pill.text_frame
        tf_pill.word_wrap = False
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = category_tag.upper()
        p_pill.alignment = PP_ALIGN.CENTER
        p_pill.font.name = FONT_TITLE
        p_pill.font.size = Pt(10)
        p_pill.font.bold = True
        p_pill.font.color.rgb = COLOR_WHITE

        # Main Title Frame
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = title_text
        p0.font.name = FONT_TITLE
        p0.font.size = Pt(22)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_HEADER_BG

        p1 = tf.add_paragraph()
        p1.text = subtitle_text
        p1.font.name = FONT_BODY
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_TEXT_MUTED
        p1.space_before = Pt(3)

        # Footer Line & Text
        tb_foot = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(10.0), Inches(0.35))
        tf_foot = tb_foot.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "SIMD-Accelerated Post-Quantum Cryptography (PQC) Framework & Benchmark Suite"
        p_foot.font.name = FONT_BODY
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = COLOR_TEXT_MUTED

        # Slide Number Badge
        tb_num = slide.shapes.add_textbox(Inches(12.0), Inches(7.0), Inches(0.6), Inches(0.35))
        tf_num = tb_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.alignment = PP_ALIGN.RIGHT
        p_num.font.name = FONT_TITLE
        p_num.font.size = Pt(12)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_PILL_BG

    def add_card(slide, left, top, width, height, title, content_list, title_color=COLOR_HEADER_BG, dark_mode=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        if dark_mode:
            shape.fill.fore_color.rgb = COLOR_HEADER_BG
            shape.line.color.rgb = COLOR_PILL_BG
        else:
            shape.fill.fore_color.rgb = COLOR_CARD_BG
            shape.line.color.rgb = COLOR_CARD_BORDER
        shape.line.width = Pt(1)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        first_item = True
        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(15)
            p_title.font.bold = True
            p_title.font.color.rgb = title_color
            p_title.space_after = Pt(6)
            first_item = False

        for item in content_list:
            if first_item:
                p = tf.paragraphs[0]
                first_item = False
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE if dark_mode else COLOR_TEXT_MAIN
            p.space_after = Pt(4)

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title Slide (Reference ODP Style)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Top Pill
    pill1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.4))
    pill1.fill.solid()
    pill1.fill.fore_color.rgb = COLOR_PILL_BG
    pill1.line.color.rgb = COLOR_PILL_BG
    p = pill1.text_frame.paragraphs[0]
    p.text = "PQC RESEARCH & BENCHMARK"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Title Card (Dark Space Navy)
    shape_t = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.4), Inches(11.333), Inches(2.2))
    shape_t.fill.solid()
    shape_t.fill.fore_color.rgb = COLOR_HEADER_BG
    shape_t.line.color.rgb = COLOR_ACCENT_TEAL
    shape_t.line.width = Pt(1.5)

    tb_t = slide1.shapes.add_textbox(Inches(1.2), Inches(1.55), Inches(10.9), Inches(1.9))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    p = tf_t.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework"
    p.font.name = FONT_TITLE
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf_t.add_paragraph()
    p.text = "High-Performance Secure Communications across Heterogeneous Architectures (x86_64 Server ↔ ARM Clients)"
    p.font.name = FONT_BODY
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(186, 230, 253)
    p.space_before = Pt(8)

    # Metadata Grid (3 Cards)
    add_card(slide1, 1.0, 3.8, 3.6, 2.8, "👤 Presenter & Team", [
        "Presented by: Systems Intern",
        "Organization: Systems & PQC Lab",
        "Team ID: PQC_Bench_2026",
        "Focus: Network Security & SIMD"
    ], title_color=COLOR_ACCENT_TEAL)

    add_card(slide1, 4.866, 3.8, 3.6, 2.8, "⚙️ Technology Stack", [
        "Core: C11 POSIX Network Stack",
        "Async I/O: Linux io_uring Engine",
        "Cryptoprimitives: liboqs & OpenSSL",
        "Formal Proofs: ProVerif (Applied Pi)"
    ], title_color=COLOR_PILL_BG)

    add_card(slide1, 8.733, 3.8, 3.6, 2.8, "🏆 Core Benchmarks", [
        "0.85 ms Single LAN Handshake",
        "850 Handshakes/sec Server Scale",
        "3.2x SIMD Vector Speedup",
        "20 Formally Verified Proofs"
    ], title_color=COLOR_EMERALD)

    # Footer
    tb_foot1 = slide1.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.333), Inches(0.4))
    p = tb_foot1.text_frame.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework & Benchmark Suite  |  Slide 1"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(slide1, "Good morning/afternoon, evaluation panel. Today I am excited to present my project: 'SIMD-Accelerated Post-Quantum Cryptography Framework and Benchmark Suite'. With quantum computing advancing, classical public-key cryptography like RSA and ECC will be broken by Shor's algorithm. Our work implements and benchmarks a production-grade C11 post-quantum secure communication stack designed for heterogeneous x86-to-ARM architectures.")

    # ==========================================
    # SLIDE 2: SOFTWARE DESIGN — Problem Statement & Objectives
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "SOFTWARE DESIGN", "Problem Statement & Research Objectives", "Addressing quantum threats in real-world heterogeneous network environments.", 2)

    # Problem Statement Card (Dark Accent Card)
    add_card(slide2, 0.8, 1.8, 11.7, 1.8, "PROBLEM STATEMENT", [
        "Design a quantum-resistant C11 secure communication framework capable of executing sub-millisecond 4-way handshakes and high-throughput authenticated data transport between server-grade x86_64 hardware and low-power embedded ARM nodes (Raspberry Pi), while guaranteeing provable formal security."
    ], title_color=COLOR_WHITE, dark_mode=True)

    # 4 Objective Cards
    add_card(slide2, 0.8, 3.8, 2.7, 3.0, "1. PQC Protocol", [
        "Construct lightweight 4-way UDP PQC handshake protocol.",
        "Combines ML-KEM-768, ML-DSA-65 & ChaCha20 AEAD."
    ], title_color=COLOR_PILL_BG)

    add_card(slide2, 3.8, 3.8, 2.7, 3.0, "2. SIMD Vector Accel", [
        "Implement dual SIMD acceleration pipeline.",
        "AVX2/BMI2 on x86 server & NEON 128-bit on ARM nodes."
    ], title_color=COLOR_ACCENT_TEAL)

    add_card(slide2, 6.8, 3.8, 2.7, 3.0, "3. Async Server Engine", [
        "Utilize Linux io_uring for non-blocking network I/O.",
        "Decouple UI logging via UNIX domain sockets."
    ], title_color=COLOR_HEADER_BG)

    add_card(slide2, 9.8, 3.8, 2.7, 3.0, "4. Formal Proofs", [
        "Build Applied Pi-Calculus model in ProVerif.",
        "Mathematically verify 20 formal security properties."
    ], title_color=COLOR_EMERALD)

    add_notes(slide2, "On Slide 2, we define our problem statement and four core objectives: 1) Designing a sub-millisecond PQC UDP protocol, 2) Optimizing lattice math with SIMD vectorization, 3) Developing a high-concurrency io_uring server decoupled via UNIX IPC, and 4) Formally verifying security properties using ProVerif.")

    # ==========================================
    # SLIDE 3: METHODOLOGY — Research Pipeline
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "METHODOLOGY", "Research & Implementation Pipeline", "A five-phase engineering methodology from primitive selection to validated execution.", 3)

    phases = [
        ("Phase 1", "Primitive Selection", "FIPS 203 ML-KEM-768 & FIPS 204 ML-DSA-65 standard parameters.", COLOR_PILL_BG),
        ("Phase 2", "Formal Verification", "ProVerif Applied Pi-Calculus modeling for Secrecy & PFS.", COLOR_ACCENT_TEAL),
        ("Phase 3", "C11 Protocol Engine", "Custom 4-way UDP handshake & ChaCha20-Poly1305 AEAD.", COLOR_HEADER_BG),
        ("Phase 4", "SIMD Acceleration", "AVX2 (x86) and NEON (ARM) polynomial vector math optimization.", COLOR_AMBER),
        ("Phase 5", "Heterogeneous Bench", "Profiling via bench_pqc, bench_nclient, and bench_aead tools.", COLOR_EMERALD)
    ]

    for idx, (ph_tag, ph_title, ph_desc, col) in enumerate(phases):
        top_y = 1.8 + (idx * 0.98)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_y), Inches(11.7), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD_BG
        shape.line.color.rgb = COLOR_CARD_BORDER
        shape.line.width = Pt(1)

        # Tag box
        tag_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(top_y + 0.15), Inches(1.5), Inches(0.55))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = col
        tag_box.line.color.rgb = col
        p = tag_box.text_frame.paragraphs[0]
        p.text = ph_tag
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        # Text Frame
        tb = slide3.shapes.add_textbox(Inches(2.6), Inches(top_y + 0.1), Inches(9.7), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = f"{ph_title} — {ph_desc}"
        p0.font.name = FONT_BODY
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TEXT_MAIN

    add_notes(slide3, "Slide 3 outlines our five-phase research methodology. We progress from standard primitive selection to formal modeling in ProVerif, C11 protocol implementation, SIMD vector acceleration, and heterogeneous hardware benchmarking across x86 and Raspberry Pi testbeds.")

    # ==========================================
    # SLIDE 4: DESIGN — Unified System Architecture & IPC
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "DESIGN", "Unified System Architecture & IPC Engine", "4-Way UDP PQC Handshake Protocol and POSIX IPC Terminal Isolation Daemon.", 4)

    add_card(slide4, 0.8, 1.8, 5.7, 5.0, "🔄 4-Way PQC Handshake Flow", [
        "1. PKT_CLIENT_HELLO:",
        "   Client ──► Server: (KEM_PK_c, DSA_PK_c, Nonce_c)",
        "",
        "2. PKT_SERVER_HELLO:",
        "   Server ──► Client: (DSA_PK_s, KEM_CT, SessionID, Nonce_s, Sig_s)",
        "",
        "3. PKT_CLIENT_AUTH:",
        "   Client ──► Server: (Sig_c over full transcript)",
        "",
        "4. PKT_ACK & Active Session:",
        "   Server ──► Client: ACK packet. Derive key via SHAKE256.",
        "   Symmetric AEAD active (ChaCha20-Poly1305)."
    ], title_color=COLOR_PILL_BG)

    add_card(slide4, 6.8, 1.8, 5.7, 5.0, "🖥️ Decoupled POSIX IPC Server", [
        "• Server Core Daemon (server):",
        "   Handles high-speed UDP packet routing on port 9877 and io_uring crypto scheduling.",
        "",
        "• Receiver Terminal (server_rx):",
        "   Connects via /tmp/pqc_rx.sock UNIX Domain Socket. Pushes decrypted text & metrics.",
        "",
        "• Sender Console (server_tx):",
        "   Connects via /tmp/pqc_tx.sock UNIX Socket. Enables operator message transmission.",
        "",
        "• Zero Display Overhead:",
        "   Decouples UI rendering from network routing."
    ], title_color=COLOR_HEADER_BG)

    add_notes(slide4, "Slide 4 presents our protocol sequence and server architecture. The 4-way UDP handshake exchanges KEM keys, ciphertexts, nonces, and transcript signatures. To ensure zero logging delay during packet routing, the server daemon communicates with receiver and sender display terminals via UNIX domain sockets.")

    # ==========================================
    # SLIDE 5: FORMULATION — Cryptographic Primitives
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "CRYPTOGRAPHY", "Cryptographic Specifications & Primitives", "NIST Category 3 Standardized Cipher Suite Selection", 5)

    rows, cols = 5, 5
    table_shape5 = slide5.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.2))
    table5 = table_shape5.table

    table5.columns[0].width = Inches(1.5)
    table5.columns[1].width = Inches(2.8)
    table5.columns[2].width = Inches(2.2)
    table5.columns[3].width = Inches(2.8)
    table5.columns[4].width = Inches(2.4)

    headers5 = ["Primitive", "Standard / Algorithm", "Security Category", "Primary Function", "Key / Output Sizes"]
    data5 = [
        ["KEM", "ML-KEM-768 (Kyber768)", "NIST Cat 3 (~AES-192)", "Ephemeral Session Key Exchange", "PK: 1,184 B | CT: 1,088 B"],
        ["DSA", "ML-DSA-65 (Dilithium3)", "NIST Cat 3 (~AES-192)", "Mutual Auth & Transcript Signing", "PK: 1,952 B | Sig: 3,293 B"],
        ["AEAD", "ChaCha20-Poly1305", "256-bit Symmetric Key", "Authenticated Payload Encryption", "Key: 32 B | Nonce: 12 B | Tag: 16 B"],
        ["Hash & KDF", "SHAKE256", "Extendable Output (XOF)", "Transcript Binding & Key Derivation", "Variable Output Length"]
    ]

    for col_idx, text in enumerate(headers5):
        cell = table5.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for row_idx, row_data in enumerate(data5):
        for col_idx, text in enumerate(row_data):
            cell = table5.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_MAIN

    add_card(slide5, 0.8, 5.3, 11.7, 1.5, "💡 Hybrid Operational Security Guarantee", [
        "• Ephemeral ML-KEM-768 key encapsulation guarantees forward-secret session key derivation.",
        "• Long-term ML-DSA-65 digital signatures verify client and server identities.",
        "• SHAKE256 derives a 256-bit secret for ChaCha20-Poly1305 AEAD streaming."
    ], title_color=COLOR_ACCENT_TEAL)

    add_notes(slide5, "Slide 5 details our cryptographic specifications. We use NIST Category 3 (~AES-192 equivalent). ML-KEM-768 manages key exchange, ML-DSA-65 handles transcript signatures, and ChaCha20-Poly1305 AEAD provides symmetric encryption after SHAKE256 key derivation.")

    # ==========================================
    # SLIDE 6: FORMAL VERIFICATION — ProVerif Analysis
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "SECURITY", "Formal Security Verification (ProVerif)", "Applied Pi-Calculus mathematical proof verifying 20 distinct security properties.", 6)

    rows, cols = 6, 3
    table_shape6 = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    table6 = table_shape6.table
    table6.columns[0].width = Inches(2.4)
    table6.columns[1].width = Inches(1.5)
    table6.columns[2].width = Inches(3.6)

    headers6 = ["Security Property", "Result", "ProVerif Formal Query"]
    data6 = [
        ["Session Key Secrecy", "VERIFIED", "query attacker(session_key_secret)"],
        ["Client Payload Secrecy", "VERIFIED", "query attacker(client_data)"],
        ["Server Payload Secrecy", "VERIFIED", "query attacker(server_data)"],
        ["Injective Mutual Auth", "VERIFIED", "inj-event(server_accepts) ==> client_initiates"],
        ["Perfect Forward Secrecy", "VERIFIED", "Phase 1 long-term key leak query"]
    ]

    for col_idx, text in enumerate(headers6):
        cell = table6.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for row_idx, row_data in enumerate(data6):
        for col_idx, text in enumerate(row_data):
            cell = table6.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            if col_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_EMERALD
            else:
                p.font.color.rgb = COLOR_TEXT_MAIN

    add_card(slide6, 8.6, 1.8, 3.9, 4.8, "🔍 Verification Highlights", [
        "• Dolev-Yao Attacker Model:",
        "   Active adversary controls physical network (intercept, modify, inject, replay).",
        "",
        "• Phase-Based PFS Proof:",
        "   Phase 0: Honest sessions.",
        "   Phase 1: Long-term key leak.",
        "   Result: Phase 0 session keys remain secret.",
        "",
        "• Attack Immunity:",
        "   Resistant to KEM substitution, UKS attacks, and reflection."
    ], title_color=COLOR_EMERALD)

    add_notes(slide6, "On Slide 6, we present formal verification results using ProVerif. Modeled in Applied Pi-Calculus under the Dolev-Yao model, ProVerif verified all 20 security queries as TRUE, mathematically proving Perfect Forward Secrecy and mutual authentication.")

    # ==========================================
    # SLIDE 7: HARDWARE ACCELERATION — SIMD Optimization
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "OPTIMIZATION", "SIMD Acceleration & arm64_ref Toolchain", "Microarchitectural SIMD optimizations and isolated cross-build workspace.", 7)

    add_card(slide7, 0.8, 1.8, 5.7, 5.0, "⚡ SIMD Acceleration Pipeline", [
        "• x86_64 Server Acceleration:",
        "   AVX2 256-bit vector registers + BMI2 instruction set. Parallelizes polynomial NTT & Keccak SHAKE256 lanes.",
        "   Result: 3.2x speedup over scalar C.",
        "",
        "• ARM64 (Raspberry Pi 4/5):",
        "   128-bit NEON SIMD vector registers.",
        "   Result: 2.8x speedup over scalar ARM C.",
        "",
        "• ARM32 (Raspberry Pi 3):",
        "   NEON + VFPv4 hard-float execution.",
        "   Result: 2.1x speedup over scalar ARM32."
    ], title_color=COLOR_PILL_BG)

    add_card(slide7, 6.8, 1.8, 5.7, 5.0, "🛠️ Portable arm64_ref Workspace", [
        "• Cross-Compilation Challenge:",
        "   Compiling ARM binaries on x86 host leads to OpenSSL static library architecture mismatches (ELFCLASS64).",
        "",
        "• Standalone Reference Layout:",
        "   Engineered arm64_ref/ containing ARM64 pre-compiled static libraries (libcrypto.a, liboqs.a) and matching headers.",
        "",
        "• Immediate Cross-Build:",
        "   Enables immediate cross-compilation via aarch64-linux-gnu-gcc without system multi-arch pollution."
    ], title_color=COLOR_HEADER_BG)

    add_notes(slide7, "Slide 7 details our hardware acceleration. AVX2 on x86 servers and NEON on ARM devices accelerate Number Theoretic Transforms in lattice math by up to 3.2x. We also created the arm64_ref toolchain workspace for seamless x86-to-ARM cross-compilation.")

    # ==========================================
    # SLIDE 8: EXPERIMENTAL SETUP — Testbed & Profiling
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "BENCHMARKING", "Heterogeneous Testbed & Profiling Suite", "Multi-hardware evaluation matrix and custom benchmark executables.", 8)

    add_card(slide8, 0.8, 1.8, 2.7, 1.8, "AMD Ryzen", ["x86_64 Server", "16 GB DDR4 RAM", "Ubuntu 24.04 LTS"], title_color=COLOR_HEADER_BG)
    add_card(slide8, 3.8, 1.8, 2.7, 1.8, "Raspberry Pi 5", ["ARM Cortex-A76", "8 GB RAM", "ARM64 / NEON"], title_color=COLOR_PILL_BG)
    add_card(slide8, 6.8, 1.8, 2.7, 1.8, "Raspberry Pi 4B", ["ARM Cortex-A72", "4 GB RAM", "ARM64 / NEON"], title_color=COLOR_ACCENT_TEAL)
    add_card(slide8, 9.8, 1.8, 2.7, 1.8, "Raspberry Pi 3B", ["ARM Cortex-A53", "1 GB RAM", "ARM32 / NEON"], title_color=COLOR_EMERALD)

    add_card(slide8, 0.8, 3.9, 3.7, 2.9, "📊 bench_pqc", [
        "Microsecond primitive benchmarks across 1,000 runs.",
        "Profiles ML-KEM KeyGen/Encap/Decap and ML-DSA Sign/Verify."
    ], title_color=COLOR_PILL_BG)

    add_card(slide8, 4.8, 3.9, 3.7, 2.9, "🔥 bench_nclient", [
        "High-concurrency multi-client stress test.",
        "Simulates 1 to 1,000 concurrent client threads measuring throughput (HS/s) and P95/P99 latencies."
    ], title_color=COLOR_AMBER)

    add_card(slide8, 8.8, 3.9, 3.7, 2.9, "🔒 bench_aead", [
        "Symmetric ChaCha20-Poly1305 AEAD throughput tool.",
        "Evaluates payload encryption across sizes (64B to 16KB) with automated file encryption validation."
    ], title_color=COLOR_HEADER_BG)

    add_notes(slide8, "Slide 8 shows our experimental setup across an AMD Ryzen server and three generations of Raspberry Pi clients. We created three profiling tools: bench_pqc for microsecond primitive benchmarks, bench_nclient for 1000-client stress testing, and bench_aead for payload throughput.")

    # ==========================================
    # SLIDE 9: EMPIRICAL RESULTS — Telemetry & Scalability
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "RESULTS", "Performance Telemetry & Empirical Results", "Microsecond primitive latencies, sub-millisecond handshakes, and scaling.", 9)

    rows, cols = 6, 5
    table_shape9 = slide9.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.8))
    table9 = table_shape9.table
    table9.columns[0].width = Inches(2.3)
    table9.columns[1].width = Inches(1.3)
    table9.columns[2].width = Inches(1.3)
    table9.columns[3].width = Inches(1.3)
    table9.columns[4].width = Inches(1.3)

    headers9 = ["Operation (μs)", "x86 Server", "Pi 5", "Pi 4B", "Pi 3B"]
    data9 = [
        ["ML-KEM KeyGen", "14.2 μs", "45.8 μs", "98.4 μs", "245.1 μs"],
        ["ML-KEM Encap", "18.6 μs", "58.2 μs", "122.6 μs", "310.5 μs"],
        ["ML-KEM Decap", "15.1 μs", "49.6 μs", "104.2 μs", "268.0 μs"],
        ["ML-DSA Sign", "62.4 μs", "185.3 μs", "395.0 μs", "980.2 μs"],
        ["ML-DSA Verify", "24.8 μs", "76.1 μs", "158.4 μs", "412.6 μs"]
    ]

    for col_idx, text in enumerate(headers9):
        cell = table9.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

    for row_idx, row_data in enumerate(data9):
        for col_idx, text in enumerate(row_data):
            cell = table9.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_MAIN

    add_card(slide9, 8.6, 1.8, 3.9, 1.05, "0.85 ms", ["Single Client LAN Handshake RTT"], title_color=COLOR_PILL_BG)
    add_card(slide9, 8.6, 3.05, 3.9, 1.05, "850 HS/sec", ["Peak Server Handshake Throughput"], title_color=COLOR_EMERALD)
    add_card(slide9, 8.6, 4.3, 3.9, 1.05, "12.4 ms", ["99th Percentile Latency (500 Clients)"], title_color=COLOR_AMBER)
    add_card(slide9, 8.6, 5.55, 3.9, 1.05, "1.8 mJ", ["Energy / Handshake on Pi 5"], title_color=COLOR_HEADER_BG)

    add_notes(slide9, "Slide 9 presents empirical performance results. On Raspberry Pi 5, ML-KEM encapsulation takes 58 μs and ML-DSA verify takes 76 μs. The 4-way LAN handshake completes in 0.85 ms. Server throughput peaks at 850 handshakes/sec, maintaining a P99 latency of 12.4 ms under 500 concurrent connections.")

    # ==========================================
    # SLIDE 10: ROADMAP — Future Scope & Extensions
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "ROADMAP", "Future Enhancements & Research Roadmap", "Extending PQC protocols towards hybrid modes, hardware offloading, and microservices.", 10)

    add_card(slide10, 0.8, 1.8, 5.7, 5.0, "🚀 Stretch Goals & Roadmap", [
        "• Hybrid Transition Modes:",
        "   Combining ECDH X25519 with ML-KEM-768 for hybrid classical-quantum handshake security during industry transition.",
        "",
        "• Hardware Offloading:",
        "   Integrating FPGA co-processors and custom RISC-V PQC vector instruction extensions for low-power edge nodes.",
        "",
        "• Zero-Trust Containerization:",
        "   Packaging the PQC protocol stack for Kubernetes service mesh (e.g. Istio) microservice sidecars."
    ], title_color=COLOR_PILL_BG)

    add_card(slide10, 6.8, 1.8, 5.7, 5.0, "⚡ Implementation Highlights", [
        "• Fully Decoupled C11 Architecture:",
        "   Client and server binaries operate with zero external dynamic runtime dependencies.",
        "",
        "• Multi-Platform Target Support:",
        "   Build scripts supporting native and cross-compilation for ARM32, ARM64, and x86_64.",
        "",
        "• Automated Test & Benchmarking:",
        "   Integrated automated benchmark suite tracking primitive microbenchmarks, throughput, and memory consumption."
    ], title_color=COLOR_HEADER_BG)

    add_notes(slide10, "Slide 10 presents future research directions: 1) Hybrid classical-quantum modes for transition, 2) FPGA and RISC-V hardware offloading, and 3) Zero-trust containerization for microservices.")

    # ==========================================
    # SLIDE 11: SUMMARY — Conclusion & Key Takeaways
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "SUMMARY", "Conclusion & Key Internship Takeaways", "Summary of technical accomplishments, performance gains, and engineering skills developed.", 11)

    add_card(slide11, 0.8, 1.8, 3.7, 5.0, "1. Production C11 PQC", [
        "• Successfully developed a lightweight, low-latency 4-way UDP PQC handshake protocol.",
        "",
        "• Achieved sub-millisecond connection times (0.85 ms RTT).",
        "",
        "• Integrated ML-KEM-768, ML-DSA-65, and ChaCha20-Poly1305 AEAD."
    ], title_color=COLOR_PILL_BG)

    add_card(slide11, 4.8, 1.8, 3.7, 5.0, "2. Hardware Accel & Scale", [
        "• Applied SIMD vectorization (AVX2 & NEON) achieving up to 3.2x speedup.",
        "",
        "• Scaled server throughput to 850 handshakes/sec under 500 concurrent connections.",
        "",
        "• Reduced Pi 5 energy per handshake down to 1.8 mJ."
    ], title_color=COLOR_ACCENT_TEAL)

    add_card(slide11, 8.8, 1.8, 3.7, 5.0, "3. Formally Verified", [
        "• Built Applied Pi-Calculus model in ProVerif verifier.",
        "",
        "• Mathematically proved 20 formal security properties.",
        "",
        "• Proven Perfect Forward Secrecy (PFS), mutual auth, and replay resistance."
    ], title_color=COLOR_EMERALD)

    add_notes(slide11, "In summary: 1) We built a production C11 PQC framework achieving sub-millisecond handshakes, 2) SIMD hardware acceleration yielded 3.2x speedup scaling to 850 HS/s, and 3) ProVerif formally verified 20 security properties including Perfect Forward Secrecy.")

    # ==========================================
    # SLIDE 12: THANK YOU — Questions & Discussion
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)

    pill12 = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.2), Inches(3.7), Inches(0.45))
    pill12.fill.solid()
    pill12.fill.fore_color.rgb = COLOR_PILL_BG
    pill12.line.color.rgb = COLOR_PILL_BG
    p = pill12.text_frame.paragraphs[0]
    p.text = "THANK YOU FOR YOUR TIME"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    shape_ty = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.9), Inches(10.333), Inches(4.0))
    shape_ty.fill.solid()
    shape_ty.fill.fore_color.rgb = COLOR_HEADER_BG
    shape_ty.line.color.rgb = COLOR_ACCENT_TEAL
    shape_ty.line.width = Pt(1.5)

    tb_ty = slide12.shapes.add_textbox(Inches(1.8), Inches(2.3), Inches(9.7), Inches(3.2))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True

    p = tf_ty.paragraphs[0]
    p.text = "Questions & Discussion"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf_ty.add_paragraph()
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework & Benchmark Suite"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(186, 230, 253)
    p.space_before = Pt(16)

    p = tf_ty.add_paragraph()
    p.text = "Open Floor for Q&A, Technical Evaluation, and Discussion"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(20)

    tb_foot12 = slide12.shapes.add_textbox(Inches(1.0), Inches(6.9), Inches(11.333), Inches(0.4))
    p = tb_foot12.text_frame.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework & Benchmark Suite  |  Slide 12"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_TEXT_MUTED

    add_notes(slide12, "Thank you to my mentors and the evaluation committee! I am now happy to open the floor to any questions and technical discussion.")

    # Save files
    pptx_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.pptx"
    prs.save(pptx_path)
    print(f"Refined PPTX saved at: {pptx_path}")

    # Convert to ODP
    cmd = f"libreoffice --headless --convert-to odp '{pptx_path}' --outdir '/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/'"
    os.system(cmd)

    odp_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.odp"

    # Copy to user home directory as well for easy access
    os.system(f"cp '{odp_path}' '/home/vishal-kumar/PQC_Internship_Presentation.odp'")
    os.system(f"cp '{pptx_path}' '/home/vishal-kumar/PQC_Internship_Presentation.pptx'")
    print(f"ODP successfully generated and copied to /home/vishal-kumar/PQC_Internship_Presentation.odp")

if __name__ == "__main__":
    build_refined_pqc_presentation()
