import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_pqc_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Design Theme Color Palette
    COLOR_BG = RGBColor(11, 15, 25)         # Deep slate dark
    COLOR_CARD = RGBColor(22, 30, 46)       # Glass card dark
    COLOR_TEXT = RGBColor(243, 244, 246)    # Light gray text
    COLOR_MUTED = RGBColor(156, 163, 175)   # Muted gray text
    COLOR_CYAN = RGBColor(6, 182, 212)      # Vibrant cyan accent
    COLOR_BLUE = RGBColor(59, 130, 246)     # Vibrant blue accent
    COLOR_PURPLE = RGBColor(139, 92, 246)   # Vibrant purple accent
    COLOR_EMERALD = RGBColor(16, 185, 129)  # Success emerald green
    COLOR_AMBER = RGBColor(245, 158, 11)    # Warning amber

    FONT_TITLE = "Inter"
    FONT_BODY = "Inter"
    FONT_MONO = "DejaVu Sans Mono"

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, slide_num_str, title_text, subtitle_text):
        # Top Cyan Progress Bar
        shape_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
        shape_bar.fill.solid()
        shape_bar.fill.fore_color.rgb = COLOR_CYAN
        shape_bar.line.color.rgb = COLOR_CYAN

        # Header Container
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True

        p0 = tf.paragraphs[0]
        p0.text = f"{slide_num_str}  |  {title_text}"
        p0.font.name = FONT_TITLE
        p0.font.size = Pt(24)
        p0.font.bold = True
        p0.font.color.rgb = COLOR_TEXT

        p1 = tf.add_paragraph()
        p1.text = subtitle_text
        p1.font.name = FONT_BODY
        p1.font.size = Pt(13)
        p1.font.color.rgb = COLOR_MUTED
        p1.space_before = Pt(4)

    def add_card(slide, left, top, width, height, title, content_list, title_color=COLOR_CYAN):
        # Card Background Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = RGBColor(45, 55, 75)
        shape.line.width = Pt(1)

        # Text Frame
        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(height - 0.3))
        tf = tb.text_frame
        tf.word_wrap = True

        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = FONT_TITLE
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = title_color
            p_title.space_after = Pt(8)
            first_item = True
        else:
            first_item = False

        for item in content_list:
            if first_item and not title:
                p = tf.paragraphs[0]
                first_item = False
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(6)

    def add_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # ==========================================
    # SLIDE 1: Title & Project Overview
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)

    # Main Title Box
    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(2.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "SIMD-Accelerated Post-Quantum Cryptography Framework"
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    p = tf1.add_paragraph()
    p.text = "Heterogeneous x86_64 Server ↔ ARM Client Networking with Formal Verification & Vector Acceleration"
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT
    p.space_before = Pt(10)

    # 3 Highlight Cards
    add_card(slide1, 1.0, 3.6, 3.5, 2.2, "🔐 NIST PQC Standards", [
        "ML-KEM-768 (Kyber768) FIPS 203",
        "ML-DSA-65 (Dilithium3) FIPS 204",
        "NIST Category 3 Security (~AES-192)"
    ])
    add_card(slide1, 4.9, 3.6, 3.5, 2.2, "⚡ SIMD Vector Accel", [
        "AVX2 & BMI2 on x86_64 Server",
        "128-bit NEON & VFPv4 on ARM",
        "Up to 3.2x Latency Reduction"
    ])
    add_card(slide1, 8.8, 3.6, 3.5, 2.2, "🛡️ Formal Verification", [
        "Applied Pi-Calculus in ProVerif",
        "20 Security Properties Proven",
        "Perfect Forward Secrecy (PFS)"
    ])

    # Metadata Banner
    add_card(slide1, 1.0, 6.0, 11.333, 1.0, "📋 Internship Presentation Metadata", [
        "Presenter: Cybersecurity & Systems Engineering Intern   |   Core Tech Stack: C11, io_uring, POSIX Threads, OpenSSL 3.0, liboqs",
        "Target Platforms: x86_64 Host Server ↔ Raspberry Pi 3B / 4B / 5 Clients   |   Achievement: Sub-millisecond 4-way Handshake (0.85 ms)"
    ], title_color=COLOR_PURPLE)

    add_notes(slide1, "Good morning/afternoon, evaluation committee. Today I present my internship project: 'SIMD-Accelerated Post-Quantum Cryptography Framework and Benchmark Suite'. With quantum computing advancing, classical public-key algorithms like RSA and ECC are vulnerable to Shor's algorithm. This project implements and benchmarks a production-grade C11 post-quantum secure networking stack tailored for heterogeneous x86-to-ARM architectures.")

    # ==========================================
    # SLIDE 2: Background, Motivation & Problem Statement
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "SLIDE 02", "Background, Motivation & Problem Statement", "The Quantum Threat & Embedded Network Bottlenecks")

    add_card(slide2, 0.8, 1.8, 5.7, 5.0, "⚠️ The Quantum Threat & NIST Transition", [
        "• Shor's Algorithm Threat:",
        "   Efficiently solves integer factorization and discrete logarithms, rendering RSA-2048/4096 and ECDH/ECDSA obsolete.",
        "",
        "• 'Harvest Now, Decrypt Later':",
        "   Adversaries currently capture encrypted network traffic to decrypt once quantum hardware matures.",
        "",
        "• NIST PQC Standardization:",
        "   Mandated migration to ML-KEM (FIPS 203) for key exchange and ML-DSA (FIPS 204) for digital signatures."
    ], title_color=COLOR_AMBER)

    add_card(slide2, 6.8, 1.8, 5.7, 5.0, "🧩 Embedded & Resource Constraints", [
        "• Massive Key & Ciphertext Sizes:",
        "   ML-KEM-768 public key = 1,184 B; ML-DSA-65 signature = 3,293 B (vs classical ECC 64 B public key).",
        "",
        "• Heavy Mathematical Overhead:",
        "   Lattice polynomial vector arithmetic and Number Theoretic Transforms (NTT) strain low-power CPUs.",
        "",
        "• Embedded RAM Bottlenecks:",
        "   High parallel compilation/execution on ARM nodes risks RAM exhaustion and Out-Of-Memory (OOM) failures."
    ], title_color=COLOR_PURPLE)

    add_notes(slide2, "To understand why this work is critical: Shor's algorithm will compromise traditional public-key cryptography. NIST has standardized ML-KEM and ML-DSA. However, PQC introduces large key sizes and heavy lattice polynomial math. On resource-constrained edge platforms like Raspberry Pi, unoptimized PQC severely degrades latency and memory usage. Our work solves this via hardware acceleration.")

    # ==========================================
    # SLIDE 3: Project Scope & Key Engineering Contributions
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "SLIDE 03", "Project Scope & Key Engineering Contributions", "Four Core Engineering Pillars Implemented During Internship")

    add_card(slide3, 0.8, 1.8, 5.7, 2.4, "1. High-Performance C11 Network Engine", [
        "• Engineered a lightweight 4-way UDP PQC handshake protocol in C11.",
        "• Achieved sub-millisecond connection establishment (0.85 ms) over LAN.",
        "• Decoupled packet headers, nonces, and session key state management."
    ])

    add_card(slide3, 6.8, 1.8, 5.7, 2.4, "2. Dual SIMD Acceleration Pipeline", [
        "• Integrated 256-bit AVX2 & BMI2 instructions for x86_64 server hosts.",
        "• Harnesses 128-bit NEON vectorization and VFPv4 for ARM Cortex nodes.",
        "• Reduced lattice polynomial multiplication latency by up to 3.2x."
    ])

    add_card(slide3, 0.8, 4.5, 5.7, 2.4, "3. Decoupled POSIX IPC Server Architecture", [
        "• Utilized Linux io_uring asynchronous I/O engine for high socket concurrency.",
        "• Decoupled UDP core daemon from UI logging via UNIX domain sockets.",
        "• Integrated server_rx display and server_tx command consoles."
    ], title_color=COLOR_PURPLE)

    add_card(slide3, 6.8, 4.5, 5.7, 2.4, "4. Formally Verified Security Model", [
        "• Modeled 4-way handshake in Applied Pi-Calculus using ProVerif.",
        "• Mathematically verified 20 security properties under Dolev-Yao model.",
        "• Proven Perfect Forward Secrecy (PFS), Mutual Auth, and Replay Resistance."
    ], title_color=COLOR_EMERALD)

    add_notes(slide3, "This slide summarizes my four main contributions: 1) A C11 sub-millisecond UDP network protocol, 2) Dual SIMD hardware acceleration for x86 and ARM, 3) An asynchronous io_uring multi-client server core decoupled via IPC, and 4) Complete formal verification using ProVerif.")

    # ==========================================
    # SLIDE 4: Cryptographic Specifications & Primitives
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "SLIDE 04", "Cryptographic Specifications & Primitives", "NIST Category 3 Standardized Cipher Suite Selection")

    # Table of Cryptographic Specifications
    rows, cols = 5, 5
    table_shape = slide4.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.2))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.8)
    table.columns[4].width = Inches(2.4)

    headers = ["Primitive", "Standard / Algorithm", "Security Category", "Primary Function", "Key / Output Sizes"]
    data = [
        ["KEM", "ML-KEM-768 (Kyber768)", "NIST Cat 3 (~AES-192)", "Ephemeral Session Key Exchange", "PK: 1,184 B | CT: 1,088 B"],
        ["DSA", "ML-DSA-65 (Dilithium3)", "NIST Cat 3 (~AES-192)", "Mutual Auth & Transcript Signing", "PK: 1,952 B | Sig: 3,293 B"],
        ["AEAD", "ChaCha20-Poly1305", "256-bit Symmetric Key", "Authenticated Payload Encryption", "Key: 32 B | Nonce: 12 B | Tag: 16 B"],
        ["Hash & KDF", "SHAKE256", "Extendable Output (XOF)", "Transcript Binding & Key Derivation", "Variable Output Length"]
    ]

    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

    for row_idx, row_data in enumerate(data):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT

    add_card(slide4, 0.8, 5.3, 11.7, 1.6, "💡 Hybrid Operational Security Guarantee", [
        "• Combines forward-secret ephemeral ML-KEM-768 key exchange with ML-DSA-65 digital signature authentication.",
        "• Once the 4-way handshake establishes a 256-bit shared secret via SHAKE256 KDF, communications switch to ChaCha20-Poly1305 AEAD, ensuring high-throughput encrypted stream transport."
    ])

    add_notes(slide4, "On Slide 4, we detail our cryptographic cipher suite selection. We chose NIST Category 3 (~AES-192 equivalent). ML-KEM-768 handles ephemeral key exchange, while ML-DSA-65 provides mutual identity signing. Once established, SHAKE256 derives a 256-bit key for ChaCha20-Poly1305 AEAD streaming.")

    # ==========================================
    # SLIDE 5: Formal Security Verification (ProVerif Analysis)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "SLIDE 05", "Formal Security Verification (ProVerif Analysis)", "Applied Pi-Calculus Mathematical Proofs of Security")

    # Table of Formal Verification Properties
    rows, cols = 7, 3
    table_shape5 = slide5.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.0))
    table5 = table_shape5.table
    table5.columns[0].width = Inches(2.4)
    table5.columns[1].width = Inches(1.6)
    table5.columns[2].width = Inches(3.5)

    headers5 = ["Security Property", "Result", "ProVerif Formal Query"]
    data5 = [
        ["Session Key Secrecy", "VERIFIED", "query attacker(session_key_secret)"],
        ["Client Payload Secrecy", "VERIFIED", "query attacker(client_data)"],
        ["Server Payload Secrecy", "VERIFIED", "query attacker(server_data)"],
        ["Client Authentication", "VERIFIED", "inj-event(server_accepts) ==> client_initiates"],
        ["Server Authentication", "VERIFIED", "inj-event(client_accepts) ==> server_initiates"],
        ["Perfect Forward Secrecy", "VERIFIED", "Phase 1 long-term key leak query"]
    ]

    for col_idx, text in enumerate(headers5):
        cell = table5.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

    for row_idx, row_data in enumerate(data5):
        for col_idx, text in enumerate(row_data):
            cell = table5.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            if col_idx == 1:
                p.font.bold = True
                p.font.color.rgb = COLOR_EMERALD
            else:
                p.font.color.rgb = COLOR_TEXT

    add_card(slide5, 8.6, 1.8, 3.9, 5.0, "🔍 Verification Highlights", [
        "• Dolev-Yao Attacker Model:",
        "   Active adversary controls physical network (intercept, modify, inject, replay).",
        "",
        "• Phase-Based PFS Modeling:",
        "   Phase 0: Honest sessions.",
        "   Phase 1: Long-term signature key leaked to network.",
        "   Result: Phase 0 session keys remain completely confidential.",
        "",
        "• Attack Immunity Proved:",
        "   Mathematically proven resistant to KEM substitution, UKS attacks, and reflection."
    ], title_color=COLOR_EMERALD)

    add_notes(slide5, "Security claims require mathematical proof. On Slide 5, we show ProVerif verification results. Modeled in Applied Pi-Calculus under the Dolev-Yao model, ProVerif verified all 20 security queries as TRUE, mathematically proving Perfect Forward Secrecy and mutual authentication.")

    # ==========================================
    # SLIDE 6: 4-Way PQC Handshake & Multi-Client Architecture
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "SLIDE 06", "4-Way PQC Handshake & Multi-Client Architecture", "Protocol Message Flow & Decoupled POSIX IPC Server Engine")

    add_card(slide6, 0.8, 1.8, 6.2, 5.0, "🔄 4-Way PQC Handshake Sequence", [
        "1. PKT_CLIENT_HELLO:",
        "   Client ──► Server: (KEM_PK_c, DSA_PK_c, Nonce_c)",
        "",
        "2. PKT_SERVER_HELLO:",
        "   Server ──► Client: (DSA_PK_s, KEM_CT, SessionID, Nonce_s, Sig_s)",
        "",
        "3. PKT_CLIENT_AUTH:",
        "   Client ──► Server: (Sig_c over full transcript)",
        "",
        "4. PKT_ACK & Active Session:",
        "   Server ──► Client: ACK packet. Derive key via SHAKE256.",
        "   Symmetric AEAD active (ChaCha20-Poly1305)."
    ])

    add_card(slide6, 7.3, 1.8, 5.2, 5.0, "🖥️ Decoupled Server IPC Architecture", [
        "• Server Core Daemon (server):",
        "   Handles high-speed UDP packet routing on port 9877 and io_uring crypto scheduling.",
        "",
        "• Receiver Terminal (server_rx):",
        "   Connects via /tmp/pqc_rx.sock UNIX Domain Socket. Pushes decrypted text & metrics.",
        "",
        "• Sender Console (server_tx):",
        "   Connects via /tmp/pqc_tx.sock UNIX Socket. Enables operator message transmission.",
        "",
        "• Zero Display Latency:",
        "   Decouples UI rendering from packet routing."
    ], title_color=COLOR_PURPLE)

    add_notes(slide6, "Slide 6 illustrates our 4-way UDP handshake sequence and server IPC architecture. The handshake exchanges KEM public keys, KEM ciphertexts, nonces, and transcript signatures. The server core daemon is decoupled from display consoles via UNIX domain sockets, preventing UI logging from delaying network packet processing.")

    # ==========================================
    # SLIDE 7: Hardware Acceleration & Cross-Platform Optimization
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "SLIDE 07", "Hardware Acceleration & Cross-Platform Optimization", "SIMD Microarchitectural Acceleration & arm64_ref Toolchain")

    add_card(slide7, 0.8, 1.8, 5.7, 5.0, "⚡ SIMD Acceleration Pipeline", [
        "• x86_64 Server Acceleration:",
        "   AVX2 256-bit vector registers + BMI2 instruction set. Parallelizes polynomial NTT & Keccak SHAKE256 lanes.",
        "   Result: 3.2x speedup over scalar C.",
        "",
        "• ARM64 (Raspberry Pi 4/5):",
        "   128-bit NEON SIMD vector registers.",
        "   Result: 2.8x speedup over scalar ARM C.",
        "",
        "• ARM32 (Raspberry Pi 3):",
        "   NEON + VFPv4 hard-float execution.",
        "   Result: 2.1x speedup over scalar ARM32."
    ])

    add_card(slide7, 6.8, 1.8, 5.7, 5.0, "🛠️ Portable arm64_ref Cross-Build Toolchain", [
        "• Cross-Compilation Challenge:",
        "   Compiling ARM binaries on x86 host leads to OpenSSL static library architecture mismatches (ELFCLASS64).",
        "",
        "• Standalone Toolchain Solution:",
        "   Engineered arm64_ref/ containing ARM64 pre-compiled static libraries (libcrypto.a, liboqs.a) and matching headers.",
        "",
        "• Immediate Cross-Build:",
        "   Enables immediate cross-compilation via aarch64-linux-gnu-gcc without system multi-arch pollution."
    ], title_color=COLOR_PURPLE)

    add_notes(slide7, "Slide 7 details our hardware acceleration techniques. AVX2 on x86 servers and NEON on ARM devices accelerate Number Theoretic Transforms in lattice math by up to 3.2x. We also created the arm64_ref toolchain, allowing developers to cross-compile ARM binaries on x86 host PCs seamlessly.")

    # ==========================================
    # SLIDE 8: Experimental Setup & Benchmarking Suite
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "SLIDE 08", "Experimental Setup & Benchmarking Suite", "Heterogeneous Hardware Matrix & Custom Benchmark Executables")

    # 4 Hardware Metric Cards
    add_card(slide8, 0.8, 1.8, 2.7, 1.8, "AMD Ryzen", ["x86_64 Server", "16 GB DDR4 RAM", "Ubuntu 24.04 LTS"], title_color=COLOR_CYAN)
    add_card(slide8, 3.8, 1.8, 2.7, 1.8, "Raspberry Pi 5", ["ARM Cortex-A76", "8 GB RAM", "ARM64 / NEON"], title_color=COLOR_CYAN)
    add_card(slide8, 6.8, 1.8, 2.7, 1.8, "Raspberry Pi 4B", ["ARM Cortex-A72", "4 GB RAM", "ARM64 / NEON"], title_color=COLOR_CYAN)
    add_card(slide8, 9.8, 1.8, 2.7, 1.8, "Raspberry Pi 3B", ["ARM Cortex-A53", "1 GB RAM", "ARM32 / NEON"], title_color=COLOR_CYAN)

    # 3 Benchmark Tool Cards
    add_card(slide8, 0.8, 3.9, 3.7, 2.9, "📊 bench_pqc", [
        "Microsecond primitive benchmarks across 1,000 runs.",
        "Profiles ML-KEM KeyGen/Encap/Decap and ML-DSA Sign/Verify."
    ])

    add_card(slide8, 4.8, 3.9, 3.7, 2.9, "🔥 bench_nclient", [
        "High-concurrency multi-client stress test.",
        "Simulates 1 to 1,000 concurrent client threads measuring throughput (HS/s) and P95/P99 latencies."
    ], title_color=COLOR_AMBER)

    add_card(slide8, 8.8, 3.9, 3.7, 2.9, "🔒 bench_aead", [
        "Symmetric ChaCha20-Poly1305 AEAD throughput tool.",
        "Evaluates payload encryption across sizes (64B to 16KB) with automated file encryption validation."
    ], title_color=COLOR_PURPLE)

    add_notes(slide8, "Slide 8 presents our experimental setup across an AMD Ryzen server and three generations of Raspberry Pi clients. To evaluate performance objectively, I created three profiling tools: bench_pqc for microsecond primitive benchmarks, bench_nclient for 1000-client stress testing, and bench_aead for payload throughput.")

    # ==========================================
    # SLIDE 9: Benchmark Results & Telemetry Analysis
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "SLIDE 09", "Benchmark Results & Performance Telemetry", "Empirical Primitive Latencies & High-Concurrency Scalability")

    # Table of Cryptographic Primitive Timings
    rows, cols = 6, 5
    table_shape9 = slide9.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5.0))
    table9 = table_shape9.table
    table9.columns[0].width = Inches(2.3)
    table9.columns[1].width = Inches(1.3)
    table9.columns[2].width = Inches(1.3)
    table9.columns[3].width = Inches(1.3)
    table9.columns[4].width = Inches(1.3)

    headers9 = ["Operation (μs)", "x86 Server", "Pi 5", "Pi 4B", "Pi 3B"]
    data9 = [
        ["ML-KEM KeyGen", "14.2 μs", "45.8 μs", "98.4 μs", "245.1 μs"],
        ["ML-KEM Encap", "18.6 μs", "58.2 μs", "122.6 μs", "310.5 μs"],
        ["ML-KEM Decap", "15.1 μs", "49.6 μs", "104.2 μs", "268.0 μs"],
        ["ML-DSA Sign", "62.4 μs", "185.3 μs", "395.0 μs", "980.2 μs"],
        ["ML-DSA Verify", "24.8 μs", "76.1 μs", "158.4 μs", "412.6 μs"]
    ]

    for col_idx, text in enumerate(headers9):
        cell = table9.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_TITLE
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN

    for row_idx, row_data in enumerate(data9):
        for col_idx, text in enumerate(row_data):
            cell = table9.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT

    # 4 Performance Telemetry Stat Cards
    add_card(slide9, 8.6, 1.8, 3.9, 1.1, "0.85 ms", ["Single Client LAN Handshake RTT"], title_color=COLOR_CYAN)
    add_card(slide9, 8.6, 3.1, 3.9, 1.1, "850 HS/sec", ["Peak Server Handshake Throughput"], title_color=COLOR_EMERALD)
    add_card(slide9, 8.6, 4.4, 3.9, 1.1, "12.4 ms", ["99th Percentile Latency (500 Clients)"], title_color=COLOR_AMBER)
    add_card(slide9, 8.6, 5.7, 3.9, 1.1, "1.8 mJ", ["Energy / Handshake on Pi 5"], title_color=COLOR_PURPLE)

    add_notes(slide9, "Slide 9 presents empirical performance results. On Raspberry Pi 5, ML-KEM encapsulation takes 58 μs and ML-DSA verify takes 76 μs. The 4-way LAN handshake completes in 0.85 ms. Server throughput peaks at 850 handshakes/sec, maintaining a P99 latency of 12.4 ms under 500 concurrent connections.")

    # ==========================================
    # SLIDE 10: Internship Conclusion & Future Roadmap
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "SLIDE 10", "Internship Conclusion & Future Roadmap", "Summary of Technical Achievements & Next Steps")

    add_card(slide10, 0.8, 1.8, 5.7, 5.0, "💡 Core Technical Learnings", [
        "• Low-Level C Systems Engineering:",
        "   Asynchronous I/O using Linux io_uring, POSIX multi-threading, atomic memory orderings.",
        "",
        "• Microarchitectural Optimization:",
        "   SIMD vectorization (AVX2 & NEON) applied to lattice polynomial math.",
        "",
        "• Formal Cybersecurity Verification:",
        "   Applied Pi-Calculus protocol modeling & automated theorem proving in ProVerif.",
        "",
        "• Performance Profiling:",
        "   Telemetry profiling measuring microsecond latency, RTT, jitter, and memory allocation."
    ])

    add_card(slide10, 6.8, 1.8, 5.7, 5.0, "🚀 Future Research Roadmap", [
        "• Hybrid Transition Modes:",
        "   Combining classical ECDH X25519 with ML-KEM-768 for gradual adoption.",
        "",
        "• Hardware Offloading:",
        "   FPGA co-processor integration and custom RISC-V PQC instruction extensions.",
        "",
        "• Zero-Trust Microservices:",
        "   Packaging the PQC protocol for Kubernetes service mesh deployments.",
        "",
        "✨ Thank You! Questions & Discussion"
    ], title_color=COLOR_PURPLE)

    add_notes(slide10, "In conclusion, this internship provided invaluable experience across systems programming, SIMD acceleration, formal methods, and network security. We proved PQC can be deployed efficiently across heterogeneous networks without compromising performance. Thank you to my mentors, and I am now open for questions!")

    # Save PowerPoint file first
    pptx_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.pptx"
    prs.save(pptx_path)
    print(f"PPTX generated successfully at: {pptx_path}")

    # Convert to ODP using LibreOffice
    cmd = f"libreoffice --headless --convert-to odp '{pptx_path}' --outdir '/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/'"
    res = os.system(cmd)
    
    odp_path = "/home/vishal-kumar/Desktop/PQC_codes/Day wise work/jul 15/adaptive_benchmarking/PQC_Internship_Presentation.odp"
    if os.path.exists(odp_path):
        print(f"ODP presentation created successfully at: {odp_path}")
    else:
        print(f"LibreOffice conversion exited with status {res}")

if __name__ == "__main__":
    build_pqc_presentation()
